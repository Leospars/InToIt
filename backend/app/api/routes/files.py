"""app/api/routes/files.py — File storage and retrieval for courses using Supabase Storage"""
import io
import json
import uuid
import logging
from typing import Optional, List
from datetime import datetime

from fastapi import APIRouter, BackgroundTasks, HTTPException, UploadFile, File, Form, Header, Query
from fastapi.responses import JSONResponse, StreamingResponse
from pydantic import BaseModel
from google import genai

from app.db.supabase import get_supabase
from app.core.config import settings
from app.api.routes.ai import generate_quiz
from app.models.requests import QuizRequest

logger = logging.getLogger(__name__)

# Document parsing imports
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


router = APIRouter()

# ── Pydantic Models ────────────────────────────────────────────

class FileMetadata(BaseModel):
    id: str
    course_id: str
    filename: str
    file_type: str
    size_bytes: int
    content_type: str
    extracted_content: Optional[str] = None
    created_at: datetime
    updated_at: datetime


class FileUploadResponse(BaseModel):
    success: bool
    file_id: str
    filename: str
    extracted_content_preview: Optional[str] = None
    auto_generation_queued: bool = False
    message: str


class GeneratedTopic(BaseModel):
    topic_id: str
    name: str
    category: str
    description: Optional[str] = None
    difficulty_level: str
    prerequisites: List[str] = []
    key_concepts: List[str] = []


class FileTopicsResponse(BaseModel):
    file_id: str
    topics: List[GeneratedTopic]
    total_count: int


class FileQuizQuestion(BaseModel):
    question: str
    options: List[str]
    correct_answer: str


class FileQuiz(BaseModel):
    quiz_id: str
    file_id: str
    topic_id: str
    topic_name: str
    difficulty: str
    questions: List[FileQuizQuestion]
    created_at: datetime


class FileQuizzesResponse(BaseModel):
    file_id: str
    quizzes: List[FileQuiz]
    total_count: int


class FileListResponse(BaseModel):
    files: List[FileMetadata]
    total_count: int


class FileContentResponse(BaseModel):
    file_id: str
    filename: str
    content: str
    file_type: str
    extracted_at: datetime


# ── Helper Functions ───────────────────────────────────────────

def get_user_from_token(authorization: Optional[str]) -> str:
    """Extract and verify user from JWT token."""
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid authorization token")
    token = authorization.split(" ")[1]
    db = get_supabase()
    user = db.auth.get_user(token)
    if not user:
        raise HTTPException(status_code=401, detail="Invalid token")
    return user.user.id


def extract_pdf_content(file_bytes: bytes) -> str:
    """Extract text content from PDF file."""
    if PyPDF2 is None:
        raise HTTPException(status_code=500, detail="PDF parsing not available. Install PyPDF2.")
    
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error parsing PDF: {str(e)}")


def extract_docx_content(file_bytes: bytes) -> str:
    """Extract text content from DOCX file."""
    if DocxDocument is None:
        raise HTTPException(status_code=500, detail="DOCX parsing not available. Install python-docx.")
    
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error parsing DOCX: {str(e)}")


def extract_pptx_content(file_bytes: bytes) -> str:
    """Extract text content from PPTX file."""
    if Presentation is None:
        raise HTTPException(status_code=500, detail="PPTX parsing not available. Install python-pptx.")
    
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"\n--- Slide {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += shape.text + "\n"
            text_parts.append(slide_text)
        return "\n".join(text_parts).strip()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error parsing PPTX: {str(e)}")


def extract_txt_content(file_bytes: bytes) -> str:
    """Extract text content from TXT file."""
    try:
        return file_bytes.decode('utf-8').strip()
    except UnicodeDecodeError:
        try:
            return file_bytes.decode('latin-1').strip()
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error decoding TXT file: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error reading TXT file: {str(e)}")


def extract_file_content(file_bytes: bytes, file_type: str) -> Optional[str]:
    """Extract text content based on file type."""
    file_type = file_type.lower()
    
    if file_type == "pdf" or file_type.endswith(".pdf"):
        return extract_pdf_content(file_bytes)
    elif file_type in ["docx", "doc"] or file_type.endswith(".docx"):
        return extract_docx_content(file_bytes)
    elif file_type in ["pptx", "ppt"] or file_type.endswith(".pptx"):
        return extract_pptx_content(file_bytes)
    elif file_type == "txt" or file_type.endswith(".txt"):
        return extract_txt_content(file_bytes)
    
    return None


def get_file_extension(filename: str) -> str:
    """Get file extension from filename."""
    return filename.split(".")[-1].lower() if "." in filename else ""


# ── Auto-Generation ──────────────────────────────────────────
async def file_generate_quizzes(file_id: str, course_id: str, filename: str, extracted_content: str, num_quizzes: int = 5, user_id: Optional[str] = None):
    """Background task: use Gemini to extract topics and pre-generate quizzes from uploaded file content."""
    try:
        db = get_supabase()
        client = genai.Client(api_key=settings.gemini_api_key)
        content_snippet = extracted_content[:8000]

        # ── Step 1: Extract topics ─────────────────────────────
        topic_prompt = f"""Analyze this document and identify the key learning topics it covers.

Document: "{filename}"
Content:
{content_snippet}

Return a JSON object with this exact structure:
{{
  "topics": [
    {{
      "topic_id": "snake_case_id",
      "name": "Human Readable Name",
      "category": "category_name",
      "description": "1-2 sentence description",
      "difficulty_level": "easy",
      "prerequisites": [],
      "key_concepts": ["concept1", "concept2"]
    }}
  ]
}}

Rules:
- Generate 3-8 topics based on the document content
- topic_id must be snake_case
- prerequisites must reference other topic_ids from THIS document only
- difficulty_level must be one of: easy, medium, hard"""

        topics_response = client.models.generate_content(
            model="gemini-2.0-flash",
            contents=topic_prompt,
            config={"response_mime_type": "application/json"}
        )
        topics_data = json.loads(topics_response.text)
        topics = topics_data.get("topics", [])

        # ── Step 2: Store topics + generate quizzes using ai.py ────────────
        for topic in topics:
            namespaced_id = f"{file_id[:8]}_{topic['topic_id']}"

            try:
                db.table("topics").upsert({
                    "topic_id": namespaced_id,
                    "name": topic["name"],
                    "category": topic.get("category", "general"),
                    "description": topic.get("description"),
                    "difficulty_level": topic.get("difficulty_level", "medium"),
                    "prerequisites": [
                        f"{file_id[:8]}_{p}" for p in topic.get("prerequisites", [])
                    ],
                    "source_file_id": file_id,
                }, on_conflict="topic_id").execute()
            except Exception as e:
                logger.warning("Failed to upsert topic %s: %s", namespaced_id, e)

            # Use the proper generate_quiz function from ai.py
            try:
                quiz_request = QuizRequest(
                    topic=topic['name'],
                    difficulty=topic.get('difficulty_level', 'medium'),
                    num_questions=num_quizzes
                )
                
                # Pass profile_data as user_id to generate_quiz_service (it expects user_id for context)
                quiz_response_text = await generate_quiz(quiz_request, user_id)
                quiz_data = json.loads(quiz_response_text)

                db.table("file_quizzes").insert({
                    "file_id": file_id,
                    "topic_id": namespaced_id,
                    "topic_name": topic["name"],
                    "quiz_data": quiz_data,
                    "difficulty": topic.get("difficulty_level", "medium"),
                }).execute()
            except Exception as e:
                logger.warning("Failed to generate/store quiz for topic %s: %s", namespaced_id, e)

        logger.info("Auto-generation complete for file %s: %d topics processed", file_id, len(topics))

    except Exception as e:
        logger.error("file_generate_quizzes failed for file %s: %s", file_id, e)


# ── Endpoints ────────────────────────────────────────────────

@router.post("/upload", response_model=FileUploadResponse)
async def upload_file(
    background_tasks: BackgroundTasks,
    course_id: str = Form(...),
    file: UploadFile = File(...),
    extract_content: bool = Form(True),
    authorization: Optional[str] = Header(None)
):
    """
    Upload a file to Supabase Storage for a specific course.
    Supports PDF, DOCX, PPTX, and TXT files with optional content extraction.
    """
    user_id = get_user_from_token(authorization)
    
    # Validate course exists
    db = get_supabase()
    try:
        course_response = db.table("courses").select("course_id").eq("course_id", course_id).single().execute()
        if not course_response.data:
            raise HTTPException(status_code=404, detail="Course ID not found")
    except Exception:
        raise HTTPException(status_code=404, detail="Course ID not found")
    
    # Validate file type
    allowed_extensions = {"pdf", "docx", "pptx", "doc", "ppt", "txt"}
    file_ext = get_file_extension(file.filename)
    
    # Read file content first to get size
    file_content = await file.read()
    file_size = len(file_content)
    
    if file_ext not in allowed_extensions:
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported file type: {file_ext}. Allowed types: {', '.join(sorted(allowed_extensions))}. File '{file.filename}' cannot be processed."
        )
    
    if file_size == 0:
        raise HTTPException(status_code=400, detail=f"File '{file.filename}' is empty. Please select a file with content.")
    
    if file_size > 50 * 1024 * 1024:  # 50MB limit
        raise HTTPException(
            status_code=400, 
            detail=f"File '{file.filename}' exceeds 50MB limit. Current size: {file_size / (1024*1024):.2f}MB. Please compress the file or upload a smaller version."
        )
    
    # Generate unique file ID and storage path
    file_id = str(uuid.uuid4())
    storage_path = f"courses/{course_id}/{file_id}/{file.filename}"
    
    # Extract content if requested and file type supports it
    extracted_content = None
    content_preview = None
    
    if extract_content:
        try:
            extracted_content = extract_file_content(file_content, file_ext)
            if extracted_content:
                content_preview = extracted_content[:500] + "..." if len(extracted_content) > 500 else extracted_content
        except HTTPException:
            raise HTTPException(status_code=400, detail="Failed to extract content from file")
        except Exception:
            raise HTTPException(status_code=400, detail="Failed to extract content from file")
    
    # Upload to Supabase Storage
    db = get_supabase()
    try:
        storage_response = db.storage.from_("course-files").upload(
            path=storage_path,
            file=file_content,
            file_options={"content-type": file.content_type or "application/octet-stream"}
        )
    except Exception as e:
        # Create bucket if it doesn't exist
        try:
            db.storage.create_bucket("course-files", {"public": False})
            storage_response = db.storage.from_("course-files").upload(
                path=storage_path,
                file=file_content,
                file_options={"content-type": file.content_type or "application/octet-stream"}
            )
        except Exception as e2:
            raise HTTPException(
                status_code=500, 
                detail=f"Failed to upload file '{file.filename}' to storage. Error: {str(e2)}. Please try again later or contact support."
            )
    
    # Store metadata in database
    try:
        file_record = {
            "id": file_id,
            "course_id": course_id,
            "filename": file.filename,
            "file_type": file_ext,
            "size_bytes": file_size,
            "storage_path": storage_path,
            "extracted_content": extracted_content,
            "created_at": datetime.utcnow().isoformat(),
            "updated_at": datetime.utcnow().isoformat()
        }
        
        db.table("course_files").insert(file_record).execute()
    except Exception as e:
        # If DB insert fails, try to delete the uploaded file
        try:
            db.storage.from_("course-files").remove([storage_path])
        except Exception:
            pass
        raise HTTPException(
            status_code=500, 
            detail=f"Failed to save file metadata for '{file.filename}' to database. Error: {str(e)}. The file was uploaded but not registered. Please contact support."
        )
    
    auto_gen_queued = False
    if extracted_content:
        background_tasks.add_task(
            file_generate_quizzes, file_id, course_id, file.filename, extracted_content, 5, user_id,
        )
        auto_gen_queued = True

    return FileUploadResponse(
        success=True,
        file_id=file_id,
        filename=file.filename,
        storage_path=storage_path,
        extracted_content_preview=content_preview,
        auto_generation_queued=auto_gen_queued,
        message="File uploaded successfully. Topics and quizzes are being generated." if auto_gen_queued else "File uploaded successfully"
    )


@router.get("/{file_id}/topics", response_model=FileTopicsResponse)
async def get_file_topics(
    file_id: str,
    authorization: Optional[str] = Header(None)
):
    """Get all auto-generated topics for a specific file."""
    get_user_from_token(authorization)
    db = get_supabase()

    try:
        response = db.table("topics").select("*").eq("source_file_id", file_id).order("created_at").execute()
        topics = [
            GeneratedTopic(
                topic_id=row["topic_id"],
                name=row["name"],
                category=row["category"],
                description=row.get("description"),
                difficulty_level=row["difficulty_level"],
                prerequisites=row.get("prerequisites") or [],
                key_concepts=[],
            )
            for row in response.data
        ]
        return FileTopicsResponse(file_id=file_id, topics=topics, total_count=len(topics))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to retrieve topics: {str(e)}")


@router.get("/{file_id}/quizzes", response_model=FileQuizzesResponse)
async def get_file_quizzes(
    file_id: str,
    authorization: Optional[str] = Header(None)
):
    """Get all auto-generated quizzes for a specific file."""
    get_user_from_token(authorization)
    db = get_supabase()

    try:
        response = db.table("file_quizzes").select("*").eq("file_id", file_id).order("created_at").execute()
        quizzes = []
        for row in response.data:
            raw_questions = row["quiz_data"].get("quiz", [])
            questions = [
                FileQuizQuestion(
                    question=q["question"],
                    options=q["options"],
                    correct_answer=q["correct_answer"],
                )
                for q in raw_questions
            ]
            quizzes.append(FileQuiz(
                quiz_id=row["id"],
                file_id=row["file_id"],
                topic_id=row["topic_id"],
                topic_name=row["topic_name"],
                difficulty=row["difficulty"],
                questions=questions,
                created_at=datetime.fromisoformat(row["created_at"].replace("Z", "+00:00")),
            ))
        return FileQuizzesResponse(file_id=file_id, quizzes=quizzes, total_count=len(quizzes))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to retrieve quizzes: {str(e)}")


@router.get("/course/{course_id}", response_model=FileListResponse)
async def list_course_files(
    course_id: str,
    limit: int = Query(50, ge=1, le=100),
    offset: int = Query(0, ge=0),
    authorization: Optional[str] = Header(None)
):
    """List all files for a specific course."""
    user_id = get_user_from_token(authorization)
    
    db = get_supabase()
    
    try:
        response = db.table("course_files").select("*").eq("course_id", course_id).order("created_at", desc=True).range(offset, offset + limit - 1).execute()
        
        files = []
        for record in response.data:
            files.append(FileMetadata(
                id=record["id"],
                course_id=record["course_id"],
                filename=record["filename"],
                file_type=record["file_type"],
                size_bytes=record["size_bytes"],
                storage_path=record["storage_path"],
                extracted_content=record.get("extracted_content"),
                created_at=datetime.fromisoformat(record["created_at"].replace("Z", "+00:00")),
                updated_at=datetime.fromisoformat(record["updated_at"].replace("Z", "+00:00"))
            ))
        
        # Get total count
        count_response = db.table("course_files").select("id", count="exact").eq("course_id", course_id).execute()
        total_count = count_response.count if hasattr(count_response, 'count') else len(files)
        
        return FileListResponse(files=files, total_count=total_count)
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to list files: {str(e)}")


@router.get("/{file_id}/content", response_model=FileContentResponse)
async def get_file_content(
    file_id: str,
    authorization: Optional[str] = Header(None)
):
    """
    Retrieve and extract content from a stored file in database.
    For PDF, DOCX, and PPTX files, returns extracted text content.
    """
    user_id = get_user_from_token(authorization)
    
    db = get_supabase()
    
    # Get file metadata
    try:
        response = db.table("course_files").select("*").eq("id", file_id).single().execute()
        if not response.data:
            raise HTTPException(status_code=404, detail="File not found")
        file_record = response.data
    except Exception as e:
        raise HTTPException(status_code=404, detail="File not found")
    
    # If content was already extracted and stored, return it
    if file_record.get("extracted_content"):
        return FileContentResponse(
            file_id=file_id,
            filename=file_record["filename"],
            content=file_record["extracted_content"],
            file_type=file_record["file_type"],
            extracted_at=datetime.fromisoformat(file_record["updated_at"].replace("Z", "+00:00"))
        )
    
    # Otherwise, extract content from file bytes in database
    file_ext = file_record["file_type"]
    
    # Get file content from database
    file_bytes = file_record["file_content"]
    if isinstance(file_bytes, str):
        # Handle case where content might be base64 encoded
        import base64
        try:
            file_bytes = base64.b64decode(file_bytes)
        except Exception:
            # If decoding fails, treat as string bytes
            file_bytes = file_bytes.encode('latin-1')
    
    # Extract content
    try:
        content = extract_file_content(file_bytes, file_ext)
        if content is None:
            raise HTTPException(status_code=400, detail=f"Content extraction not supported for {file_ext} files")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Content extraction failed: {str(e)}")
    
    # Update the database with extracted content
    try:
        db.table("course_files").update({
            "extracted_content": content,
            "updated_at": datetime.utcnow().isoformat()
        }).eq("id", file_id).execute()
    except Exception:
        # Non-critical error, continue without updating
        pass
    
    return FileContentResponse(
        file_id=file_id,
        filename=file_record["filename"],
        content=content,
        file_type=file_record["file_type"],
        extracted_at=datetime.utcnow()
    )


@router.get("/{file_id}/download")
async def download_file(
    file_id: str,
    authorization: Optional[str] = Header(None)
):
    """Download a file directly from database."""
    user_id = get_user_from_token(authorization)
    
    db = get_supabase()
    
    # Get file metadata and content
    try:
        response = db.table("course_files").select("*").eq("id", file_id).single().execute()
        if not response.data:
            raise HTTPException(status_code=404, detail="File not found")
        file_record = response.data
    except Exception as e:
        raise HTTPException(status_code=404, detail="File not found")
    
    # Get file content from database
    file_bytes = file_record["file_content"]
    if isinstance(file_bytes, str):
        # Handle case where content might be base64 encoded
        import base64
        try:
            file_bytes = base64.b64decode(file_bytes)
        except Exception:
            # If decoding fails, treat as string bytes
            file_bytes = file_bytes.encode('latin-1')
    
    # Determine content type
    content_type_map = {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "doc": "application/msword",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "ppt": "application/vnd.ms-powerpoint"
    }
    content_type = file_record.get("content_type") or content_type_map.get(file_record["file_type"], "application/octet-stream")
    
    return StreamingResponse(
        io.BytesIO(file_bytes),
        media_type=content_type,
        headers={"Content-Disposition": f'attachment; filename="{file_record["filename"]}"'}
    )


@router.delete("/{file_id}")
async def delete_file(
    file_id: str,
    authorization: Optional[str] = Header(None)
):
    """Delete a file from database."""
    user_id = get_user_from_token(authorization)
    
    db = get_supabase()
    
    # Get file metadata for response message
    try:
        response = db.table("course_files").select("*").eq("id", file_id).single().execute()
        if not response.data:
            raise HTTPException(status_code=404, detail="File not found")
        file_record = response.data
    except Exception as e:
        raise HTTPException(status_code=404, detail="File not found")
    
    # Delete from database (this removes both metadata and file content)
    try:
        db.table("course_files").delete().eq("id", file_id).execute()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to delete file record: {str(e)}")
    
    return JSONResponse(content={
        "success": True,
        "message": f"File '{file_record['filename']}' deleted successfully",
        "file_id": file_id
    })


@router.get("/{file_id}")
async def get_file_metadata(
    file_id: str,
    authorization: Optional[str] = Header(None)
):
    """Get metadata for a specific file."""
    user_id = get_user_from_token(authorization)
    
    db = get_supabase()
    
    try:
        response = db.table("course_files").select("*").eq("id", file_id).single().execute()
        if not response.data:
            raise HTTPException(status_code=404, detail="File not found")
        file_record = response.data
        
        return FileMetadata(
            id=file_record["id"],
            course_id=file_record["course_id"],
            filename=file_record["filename"],
            file_type=file_record["file_type"],
            size_bytes=file_record["size_bytes"],
            storage_path=file_record["storage_path"],
            extracted_content=file_record.get("extracted_content"),
            created_at=datetime.fromisoformat(file_record["created_at"].replace("Z", "+00:00")),
            updated_at=datetime.fromisoformat(file_record["updated_at"].replace("Z", "+00:00"))
        )
    except Exception as e:
        if isinstance(e, HTTPException):
            raise
        raise HTTPException(status_code=404, detail="File not found")


@router.post("/{file_id}/reextract")
async def reextract_content(
    file_id: str,
    authorization: Optional[str] = Header(None)
):
    """Re-extract content from a stored file in database (useful if extraction failed previously)."""
    user_id = get_user_from_token(authorization)
    
    db = get_supabase()
    
    # Get file metadata
    try:
        response = db.table("course_files").select("*").eq("id", file_id).single().execute()
        if not response.data:
            raise HTTPException(status_code=404, detail="File not found")
        file_record = response.data
    except Exception as e:
        raise HTTPException(status_code=404, detail="File not found")
    
    file_ext = file_record["file_type"]
    
    # Get file content from database
    file_bytes = file_record["file_content"]
    if isinstance(file_bytes, str):
        # Handle case where content might be base64 encoded
        import base64
        try:
            file_bytes = base64.b64decode(file_bytes)
        except Exception:
            # If decoding fails, treat as string bytes
            file_bytes = file_bytes.encode('latin-1')
    
    # Extract content
    try:
        content = extract_file_content(file_bytes, file_ext)
        if content is None:
            raise HTTPException(status_code=400, detail=f"Content extraction not supported for {file_ext} files")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Content extraction failed: {str(e)}")
    
    # Update database
    try:
        db.table("course_files").update({
            "extracted_content": content,
            "updated_at": datetime.utcnow().isoformat()
        }).eq("id", file_id).execute()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to update file record: {str(e)}")
    
    return FileContentResponse(
        file_id=file_id,
        filename=file_record["filename"],
        content=content,
        file_type=file_record["file_type"],
        extracted_at=datetime.utcnow()
    )
